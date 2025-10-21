from __future__ import annotations

from pathlib import Path

from deckdown.cli import EXIT_OK, main


def test_assemble_table_and_basic(tmp_path: Path) -> None:
    # Minimal deck with a table and a rectangle
    md = tmp_path / "deck.md"
    md.write_text(
        "# t\n\n---\n```json\n{\n  \"version\": \"deckdown-1\",\n  \"slide\": {\n    \"index\": 1,\n    \"size\": {\"width_emu\": 9144000, \"height_emu\": 5143500},\n    \"shapes\": [\n      {\n        \"id\": \"s1\",\n        \"kind\": \"table\",\n        \"bbox\": {\"x_emu\": 914400, \"y_emu\": 914400, \"w_emu\": 1828800, \"h_emu\": 914400, \"x_norm\": 0.1, \"y_norm\": 0.1, \"w_norm\": 0.2, \"h_norm\": 0.178},\n        \"z\": 0,\n        \"table\": {\n          \"rows\": 2, \"cols\": 2,\n          \"cells\": [\n            { \"r\": 0, \"c\": 0, \"rowspan\": 1, \"colspan\": 1, \"text\": { \"paras\": [ { \"lvl\": 0, \"runs\": [ { \"text\": \"H1\" } ] } ] } },\n            { \"r\": 0, \"c\": 1, \"rowspan\": 1, \"colspan\": 1, \"text\": { \"paras\": [ { \"lvl\": 0, \"runs\": [ { \"text\": \"H2\" } ] } ] } },\n            { \"r\": 1, \"c\": 0, \"rowspan\": 1, \"colspan\": 1, \"text\": { \"paras\": [ { \"lvl\": 0, \"runs\": [ { \"text\": \"A\" } ] } ] } },\n            { \"r\": 1, \"c\": 1, \"rowspan\": 1, \"colspan\": 1, \"text\": { \"paras\": [ { \"lvl\": 0, \"runs\": [ { \"text\": \"B\" } ] } ] } }\n          ]\n        }\n      },\n      {\n        \"id\": \"s2\",\n        \"kind\": \"shape_basic\",\n        \"geom\": \"rectangle\",\n        \"bbox\": {\"x_emu\": 3657600, \"y_emu\": 914400, \"w_emu\": 914400, \"h_emu\": 914400, \"x_norm\": 0.4, \"y_norm\": 0.1, \"w_norm\": 0.1, \"h_norm\": 0.178},\n        \"z\": 1\n      }\n    ]\n  }\n}\n```\n",
        encoding="utf-8",
    )
    out = tmp_path / "out.pptx"
    code = main(["assemble", str(md), "-o", str(out)])
    assert code == EXIT_OK
    assert out.exists()
    from pptx import Presentation

    prs = Presentation(str(out))
    s = prs.slides[0]
    # Ensure at least one table and one auto shape exist
    types = [sh.shape_type for sh in s.shapes]
    assert any(str(t).startswith("TABLE") for t in types)
    assert any(str(t).startswith("AUTO_SHAPE") for t in types)

