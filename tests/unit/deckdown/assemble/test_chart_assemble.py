from __future__ import annotations

from pathlib import Path

from deckdown.cli import EXIT_OK, main


def test_assemble_chart_column(tmp_path: Path) -> None:
    md = tmp_path / "deck.md"
    md.write_text(
        '# t\n\n---\n```json\n{\n  "version": "deckdown-1",\n  "slide": {\n    "index": 1,\n    "size": {"width_emu": 9144000, "height_emu": 5143500},\n    "shapes": [\n      {\n        "id": "s1",\n        "kind": "chart",\n        "bbox": {"x_emu": 914400, "y_emu": 914400, "w_emu": 3657600, "h_emu": 2438400, "x_norm": 0.1, "y_norm": 0.1, "w_norm": 0.4, "h_norm": 0.474},\n        "z": 0,\n        "chart": {\n          "type": "column",\n          "categories": ["A", "B", "C"],\n          "series": [\n            { "name": "S1", "values": [1,2,3] },\n            { "name": "S2", "values": [4,5,6] }\n          ],\n          "plot_area": { "has_legend": true, "legend_pos": "right" }\n        }\n      }\n    ]\n  }\n}\n```\n',
        encoding="utf-8",
    )
    out = tmp_path / "out.pptx"
    code = main(["assemble", str(md), "-o", str(out)])
    assert code == EXIT_OK

    from pptx import Presentation

    prs = Presentation(str(out))
    s = prs.slides[0]
    chart_shapes = [sh for sh in s.shapes if hasattr(sh, "has_chart") and sh.has_chart]
    assert len(chart_shapes) == 1
    ch = chart_shapes[0].chart
    # Verify counts
    plots = list(ch.plots)
    assert plots and len(plots[0].series) == 2
