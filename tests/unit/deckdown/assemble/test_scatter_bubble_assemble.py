from __future__ import annotations

from pathlib import Path

from deckdown.cli import EXIT_OK, main


def test_assemble_scatter_and_bubble(tmp_path: Path) -> None:
    md = tmp_path / "deck.md"
    md.write_text(
        '# t\n\n---\n```json\n{\n  "version": "deckdown-1",\n  "slide": {\n    "index": 1,\n    "size": {"width_emu": 9144000, "height_emu": 5143500},\n    "shapes": [\n      {\n        "id": "s1",\n        "kind": "chart",\n        "bbox": {"x_emu": 914400, "y_emu": 914400, "w_emu": 3657600, "h_emu": 2438400, "x_norm": 0.1, "y_norm": 0.1, "w_norm": 0.4, "h_norm": 0.474},\n        "z": 0,\n        "chart": {\n          "type": "scatter",\n          "categories": [],\n          "series": [ { "name": "S", "x_values": [1,2,3], "values": [4,5,6] } ]\n        }\n      }\n    ]\n  }\n}\n```\n---\n```json\n{\n  "version": "deckdown-1",\n  "slide": {\n    "index": 2,\n    "size": {"width_emu": 9144000, "height_emu": 5143500},\n    "shapes": [\n      {\n        "id": "s2",\n        "kind": "chart",\n        "bbox": {"x_emu": 914400, "y_emu": 914400, "w_emu": 3657600, "h_emu": 2438400, "x_norm": 0.1, "y_norm": 0.1, "w_norm": 0.4, "h_norm": 0.474},\n        "z": 0,\n        "chart": {\n          "type": "bubble",\n          "series": [ { "name": "B", "x_values": [1,2,3], "values": [3,2,1], "sizes": [10,20,30] } ]\n        }\n      }\n    ]\n  }\n}\n```\n',
        encoding="utf-8",
    )
    out = tmp_path / "out.pptx"
    code = main(["assemble", str(md), "-o", str(out)])
    assert code == EXIT_OK
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation(str(out))
    ch1 = prs.slides[0].shapes[0].chart
    ch2 = prs.slides[1].shapes[0].chart
    assert ch1.chart_type == XL_CHART_TYPE.XY_SCATTER
    assert ch2.chart_type == XL_CHART_TYPE.BUBBLE
