from __future__ import annotations

from pathlib import Path

from deckdown.cli import EXIT_OK, main


def test_assemble_applies_value_axis(tmp_path: Path) -> None:
    md = tmp_path / "deck.md"
    md.write_text(
        "# t\n\n---\n```json\n{\n  \"version\": \"deckdown-1\",\n  \"slide\": {\n    \"index\": 1,\n    \"size\": {\"width_emu\": 9144000, \"height_emu\": 5143500},\n    \"shapes\": [\n      {\n        \"id\": \"s1\",\n        \"kind\": \"chart\",\n        \"bbox\": {\"x_emu\": 914400, \"y_emu\": 914400, \"w_emu\": 3657600, \"h_emu\": 2438400, \"x_norm\": 0.1, \"y_norm\": 0.1, \"w_norm\": 0.4, \"h_norm\": 0.474},\n        \"z\": 0,\n        \"chart\": {\n          \"type\": \"column\",\n          \"categories\": [\"A\", \"B\", \"C\"],\n          \"series\": [ { \"name\": \"S1\", \"values\": [1,2,3] } ],\n          \"axes\": { \n            \"category\": { \"title\": \"Cats\" },\n            \"value\": { \"title\": \"Vals\", \"min\": 0.0, \"max\": 10.0, \"major_unit\": 2.0, \"format_code\": \"0.0\" }\n          }\n        }\n      }\n    ]\n  }\n}\n```\n",
        encoding="utf-8",
    )
    out = tmp_path / "out.pptx"
    code = main(["assemble", str(md), "-o", str(out)])
    assert code == EXIT_OK
    from pptx import Presentation
    prs = Presentation(str(out))
    ch = prs.slides[0].shapes[0].chart
    # Sanity: chart exists and has two axes objects for column chart
    assert getattr(ch, "category_axis", None) is not None
    assert getattr(ch, "value_axis", None) is not None
