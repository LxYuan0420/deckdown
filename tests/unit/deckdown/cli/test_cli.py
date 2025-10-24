from __future__ import annotations

import base64
from pathlib import Path

from deckdown.cli import EXIT_INPUT_ERROR, EXIT_OK, main


class TestCLIExtract:
    SAMPLE_PNG = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAASsJTYQAAAAASUVORK5CYII="
    )

    @staticmethod
    def _write_min_pptx(path: Path) -> None:
        # Arrange: create a minimal valid PPTX (no slides)
        from pptx import Presentation

        prs = Presentation()
        prs.save(str(path))

    def _write_picture_pptx(self, path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        image_path = path.with_suffix(".png")
        image_path.write_bytes(base64.b64decode(self.SAMPLE_PNG))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), Inches(2), Inches(2))
        prs.save(str(path))

    def test_writes_markdown(self, tmp_path: Path) -> None:
        # Arrange
        pptx = tmp_path / "sample.pptx"
        self._write_min_pptx(pptx)
        out = tmp_path / "out.md"
        # Act
        code = main(["extract", str(pptx), "--md-out", str(out)])
        # Assert
        assert code == EXIT_OK
        text = out.read_text(encoding="utf-8")
        assert text.startswith("# sample\n")

    def test_default_output_name(self, tmp_path: Path) -> None:
        # Arrange
        pptx = tmp_path / "Deck Name.PPTX"  # ensure case-insensitivity
        self._write_min_pptx(pptx)
        # Act
        code = main(["extract", str(pptx)])
        # Assert
        assert code == EXIT_OK
        out = tmp_path / "Deck Name.md"
        assert out.exists()

    def test_missing_input_returns_error(self, tmp_path: Path) -> None:
        # Arrange
        missing = tmp_path / "missing.pptx"
        # Act
        code = main(["extract", str(missing)])
        # Assert
        assert code == EXIT_INPUT_ERROR

    def test_input_directory_is_error(self, tmp_path: Path) -> None:
        # Arrange & Act
        code = main(["extract", str(tmp_path)])
        # Assert
        assert code == EXIT_INPUT_ERROR

    def test_md_out_directory(self, tmp_path: Path) -> None:
        # Arrange
        pptx = tmp_path / "deck.pptx"
        self._write_min_pptx(pptx)
        outdir = tmp_path / "out"
        # Act
        code = main(["extract", str(pptx), "--md-out", str(outdir)])
        # Assert
        assert code == EXIT_OK
        expected = outdir / "deck.md"
        assert expected.exists()

    def test_embed_media_base64_is_default(self, tmp_path: Path) -> None:
        pptx = tmp_path / "picture.pptx"
        self._write_picture_pptx(pptx)
        out = tmp_path / "deck.md"

        code = main(["extract", str(pptx), "--md-out", str(out)])

        assert code == EXIT_OK
        text = out.read_text(encoding="utf-8")
        assert '"data_url": "data:image/png;base64,' in text
        assert '"ref": null' in text

    def test_embed_media_refs_writes_files(self, tmp_path: Path) -> None:
        pptx = tmp_path / "picture.pptx"
        self._write_picture_pptx(pptx)
        outdir = tmp_path / "output"
        out = outdir / "deck.md"

        code = main(["extract", str(pptx), "--md-out", str(out), "--embed-media", "refs"])

        assert code == EXIT_OK
        assets_dir = outdir / "deck_assets"
        assert assets_dir.exists()
        files = list(assets_dir.iterdir())
        assert len(files) == 1
        assert files[0].suffix == ".png"
        text = out.read_text(encoding="utf-8")
        assert '"ref": "deck_assets/' in text
        assert '"data_url": null' in text
