from __future__ import annotations

from pathlib import Path

from deckdown.extractors.ast import AstExtractor
from deckdown.loader import Loader


def _make_chart_axes_points(tmp: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.dml.color import RGBColor

    p = tmp / "cap.pptx"
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    data.add_series("S1", (1, 2, 3))
    ch = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), data
    ).chart
    # value axis settings
    va = ch.value_axis
    va.minimum_scale = 0.0
    va.maximum_scale = 10.0
    va.major_unit = 2.0
    va.tick_labels.number_format = "0.0"
    va.has_title = True
    va.axis_title.text_frame.text = "Values"
    ch.category_axis.has_title = True
    ch.category_axis.axis_title.text_frame.text = "Cats"

    # Pie with per-point colors for point-level color extraction
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    data2 = CategoryChartData()
    data2.categories = ["A", "B", "C"]
    data2.add_series("S1", (1, 2, 3))
    ch2 = s2.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(4), Inches(3), data2
    ).chart
    ser = list(ch2.plots)[0].series[0]
    for i, pt in enumerate(ser.points):
        pt.format.fill.solid()
        col = [(0xAA, 0, 0), (0, 0xAA, 0), (0, 0, 0xAA)][i]
        pt.format.fill.fore_color.rgb = RGBColor(*col)

    prs.save(str(p))
    return p


def test_chart_axes_points_extraction(tmp_path: Path) -> None:
    p = _make_chart_axes_points(tmp_path)
    prs = Loader(str(p)).presentation()
    docs = AstExtractor().extract(prs)
    # Slide 1 axes
    ch = [s for s in docs[1].slide.shapes if getattr(s, "kind", None).value == "chart"][0]
    ax = ch.chart.axes
    assert ax is not None
    assert ax.category is not None and ax.category.title == "Cats"
    v = ax.value
    assert v is not None
    assert v.title == "Values"
    assert v.min == 0.0 and v.max == 10.0 and v.major_unit == 2.0
    assert v.format_code == "0.0"
    # Slide 2 points
    ch2 = [s for s in docs[2].slide.shapes if getattr(s, "kind", None).value == "chart"][0]
    s0 = ch2.chart.series[0]
    pts = s0.points or ()
    assert len(pts) >= 3
    colors = [pts[i].color.resolved_rgb for i in range(3) if pts[i].color is not None]
    assert colors == ["#AA0000", "#00AA00", "#0000AA"]
