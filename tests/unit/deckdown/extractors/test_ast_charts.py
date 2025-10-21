from __future__ import annotations

from pathlib import Path

from deckdown.extractors.ast import AstExtractor
from deckdown.loader import Loader


def _make_chart_deck(tmp: Path) -> Path:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    p = tmp / "chart.pptx"
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    data.add_series("S1", (1, 2, 3))
    data.add_series("S2", (4, 5, 6))
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), data)
    prs.save(str(p))
    return p


def test_ast_charts_basic(tmp_path: Path) -> None:
    pptx = _make_chart_deck(tmp_path)
    prs = Loader(str(pptx)).presentation()
    docs = AstExtractor().extract(prs)
    s1 = docs[1].slide
    charts = [sh for sh in s1.shapes if getattr(sh, "kind", None).value == "chart"]
    assert len(charts) == 1
    ch = charts[0].chart
    assert ch.type in ("column", "bar", "line", "pie", "donut", "unknown")
    assert tuple(ch.categories) == ("A", "B", "C")
    assert len(ch.series) == 2
    assert tuple(ch.series[0].values) == (1, 2, 3)
    assert tuple(ch.series[1].values) == (4, 5, 6)

