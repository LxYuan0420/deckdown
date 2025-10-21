from __future__ import annotations

from pathlib import Path

from deckdown.extractors.ast import AstExtractor
from deckdown.loader import Loader


def _make_table_deck(tmp: Path) -> Path:
    from pptx import Presentation

    p = tmp / "tab.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # 2x2 table
    left = top = width = height = 914400  # 1 inch emu
    tbl_shape = s.shapes.add_table(2, 2, left, top, width, height)
    tbl = tbl_shape.table
    tbl.cell(0, 0).text = "H1"
    tbl.cell(0, 1).text = "H2"
    tbl.cell(1, 0).text = "A1"
    tbl.cell(1, 1).text = "B1"

    # Add another slide with vertical merge
    s2 = prs.slides.add_slide(blank)
    tbl_shape2 = s2.shapes.add_table(3, 2, left, top, width, height)
    tbl2 = tbl_shape2.table
    tbl2.cell(0, 0).text = "X"
    tbl2.cell(1, 0).text = "X2"
    tbl2.cell(2, 0).text = "X3"
    tbl2.cell(0, 1).text = "Y"
    tbl2.cell(1, 1).text = "Y2"
    tbl2.cell(2, 1).text = "Y3"
    # vertical merge first column: (0,0) down to (2,0)
    tbl2.cell(0, 0).merge(tbl2.cell(2, 0))

    prs.save(str(p))
    return p


def test_ast_tables(tmp_path: Path) -> None:
    pptx = _make_table_deck(tmp_path)
    prs = Loader(str(pptx)).presentation()
    docs = AstExtractor().extract(prs)

    # Slide 1: simple 2x2, expect 4 cells with rowspan/colspan=1
    s1 = docs[1].slide
    tbl_shapes = [sh for sh in s1.shapes if getattr(sh, "kind", None).value == "table"]
    assert len(tbl_shapes) == 1
    t = tbl_shapes[0]
    assert t.table.rows == 2 and t.table.cols == 2
    assert len(t.table.cells) == 4
    assert all(c.rowspan == 1 and c.colspan == 1 for c in t.table.cells)

    # Slide 2: vertical merge in col 0 spanning 3 rows → expect 3*2 - (3-1) = 4 cells emitted
    s2 = docs[2].slide
    tbl2 = [sh for sh in s2.shapes if getattr(sh, "kind", None).value == "table"][0]
    assert tbl2.table.rows == 3 and tbl2.table.cols == 2
    assert len(tbl2.table.cells) == 4
    # The top-left cell should have rowspan >= 2
    top_left = next(c for c in tbl2.table.cells if c.r == 0 and c.c == 0)
    assert top_left.rowspan >= 2
