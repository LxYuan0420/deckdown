from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path
from typing import Never

from deckdown.extractors.text import (
    ParagraphSplitter,
    ShapeOrderer,
    SlideTextExtractor,
    TextExtractor,
    TitleResolver,
)
from deckdown.loader import Loader
from deckdown.renderers.markdown import MarkdownRenderer


class TestTextExtractor:
    @staticmethod
    def _make_text_basic(tmp: Path) -> Path:
        # Arrange: tiny deck with text and bullets
        from pptx import Presentation

        out = tmp / "t.pptx"
        prs = Presentation()
        layout = prs.slide_layouts[1]

        # Slide 1
        s1 = prs.slides.add_slide(layout)
        s1.shapes.title.text = "Intro"
        tf = s1.placeholders[1].text_frame
        tf.clear()
        tf.paragraphs[0].text = "Welcome to DeckDown"
        p = tf.add_paragraph()
        p.text = "Goals: Clean Markdown from PPTX"
        p = tf.add_paragraph()
        p.text = "Item 1"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "Sub 1"
        p.level = 1

        # Slide 2
        s2 = prs.slides.add_slide(layout)
        s2.shapes.title.text = "Details"
        tf2 = s2.placeholders[1].text_frame
        tf2.clear()
        tf2.paragraphs[0].text = "Alpha"
        p = tf2.add_paragraph()
        p.text = "Beta"

        # Slide 3
        s3 = prs.slides.add_slide(layout)
        tf3 = s3.placeholders[1].text_frame
        tf3.clear()
        tf3.paragraphs[0].text = "Closing"

        prs.save(str(out))
        return out

    def test_basic(self, tmp_path: Path) -> None:
        # Arrange
        pptx = self._make_text_basic(tmp_path)
        prs = Loader(str(pptx)).presentation()
        # Act
        deck = TextExtractor().extract_deck(prs, source_path=str(pptx))
        markdown = MarkdownRenderer().render(deck)
        # Assert
        expected = (
            "# t\n\n"
            "## Slide 1 — Intro\n\n"
            "### Text\n"
            "Welcome to DeckDown\n"
            "Goals: Clean Markdown from PPTX\n\n"
            "### Bullets\n"
            "- Item 1\n"
            "  - Sub 1\n\n"
            "## Slide 2 — Details\n\n"
            "### Text\n"
            "Alpha\n"
            "Beta\n\n"
            "## Slide 3 — Untitled\n\n"
            "### Text\n"
            "Closing\n"
        )
        assert markdown == expected

    def test_skip_title_shape_only(self, tmp_path: Path) -> None:
        # Arrange: slide with only a title shape should produce no body/bullets
        from pptx import Presentation

        p = tmp_path / "title_only.pptx"
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        s = prs.slides.add_slide(title_slide_layout)
        s.shapes.title.text = "Only Title"
        prs.save(str(p))

        # Act
        deck = TextExtractor().extract_deck(prs, source_path=str(p))
        # Assert
        assert deck.slides[0].title == "Only Title"
        assert deck.slides[0].text_blocks == ()
        assert deck.slides[0].bullets == ()

    def test_title_resolver_empty_presentation(self, tmp_path: Path) -> None:
        # Arrange: no slides triggers StopIteration path in TitleResolver
        from pptx import Presentation

        p = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.save(str(p))
        # Act
        deck = TextExtractor().extract_deck(prs, source_path=str(p))
        # Assert
        assert deck.title is None

    def test_title_property_exception_handling(self) -> None:
        # Arrange: fake slide whose shapes.title raises; covers exception branch
        class _FakeShapes:
            def __iter__(self) -> Iterator[object]:  # pragma: no cover - iteration not used here
                return iter(())

            @property
            def title(self) -> Never:
                raise RuntimeError("boom")

        class _FakeSlide:
            shapes = _FakeShapes()

        class _FakePRS:
            slides = [_FakeSlide()]

        slide_extractor = SlideTextExtractor(ShapeOrderer(), ParagraphSplitter())
        # Act
        title = TitleResolver(slide_extractor).derive(_FakePRS())
        # Assert
        assert title is None
