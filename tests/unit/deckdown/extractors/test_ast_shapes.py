from __future__ import annotations

from pathlib import Path

from deckdown.extractors.ast import AstExtractor
from deckdown.loader import Loader


def _make_shapes_deck(tmp: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    p = tmp / "shapes.pptx"
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Rectangle
    s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1), Inches(1), Inches(1.5), Inches(1))
    # Line
    s.shapes.add_connector(1, Inches(1), Inches(3), Inches(3), Inches(3))
    prs.save(str(p))
    return p


def test_ast_basic_and_line(tmp_path: Path) -> None:
    pptx = _make_shapes_deck(tmp_path)
    prs = Loader(str(pptx)).presentation()
    docs = AstExtractor().extract(prs)
    s1 = docs[1].slide
    kinds = [sh.kind.value for sh in s1.shapes]
    assert "shape_basic" in kinds
    assert "line" in kinds
