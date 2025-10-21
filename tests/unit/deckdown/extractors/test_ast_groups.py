from __future__ import annotations

from pathlib import Path

from deckdown.extractors.ast import AstExtractor
from deckdown.loader import Loader


def _make_group_deck(tmp: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    p = tmp / "group.pptx"
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    oval = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(3.5), Inches(1), Inches(1), Inches(1))
    s.shapes.add_group_shape([rect, oval])
    prs.save(str(p))
    return p


def test_ast_groups(tmp_path: Path) -> None:
    pptx = _make_group_deck(tmp_path)
    prs = Loader(str(pptx)).presentation()
    docs = AstExtractor().extract(prs)
    s1 = docs[1].slide
    kinds = [sh.kind.value for sh in s1.shapes]
    assert "group" in kinds
    group = next(sh for sh in s1.shapes if getattr(sh, "kind", None).value == "group")
    # expect two children IDs
    assert len(group.children) == 2
    # children carry group id
    child_groups = [sh.group for sh in s1.shapes if getattr(sh, "kind", None).value != "group"]
    assert any(cg == group.id for cg in child_groups)
