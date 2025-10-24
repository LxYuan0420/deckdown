from __future__ import annotations

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckdown.ast import ShapeKind, SlideSize
from deckdown.color.theme import ThemeResolver
from deckdown.extractors.ast import SlideWalker
from deckdown.extractors.context import ExtractContext
from deckdown.extractors.group import GroupExtractor
from deckdown.extractors.handlers.basic_line_handler import BasicShapeHandler, LineShapeHandler
from deckdown.extractors.handlers.chart_handler import ChartShapeHandler
from deckdown.extractors.handlers.picture_handler import PictureShapeHandler
from deckdown.extractors.handlers.table_handler import TableShapeHandler
from deckdown.extractors.handlers.text_handler import TextShapeHandler


def _make_walker() -> SlideWalker:
    handlers = (
        TableShapeHandler(),
        ChartShapeHandler(),
        PictureShapeHandler(),
        LineShapeHandler(),
        BasicShapeHandler(),
        TextShapeHandler(),
    )
    return SlideWalker(handlers=handlers, group_extractor=GroupExtractor(handlers=handlers))


@pytest.fixture()
def ctx() -> ExtractContext:
    size = SlideSize(width_emu=1000, height_emu=1000)
    return ExtractContext(size=size, theme=ThemeResolver(scheme={}))


def test_slide_walker_handles_text_picture_table(ctx: ExtractContext) -> None:
    walker = _make_walker()
    shapes = [
        FakeTextShape(shape_id=1, left=100, top=50, width=300, height=100, text="Hello"),
        FakePictureShape(shape_id=2, left=600, top=100, width=200, height=200),
        FakeTableShape(
            shape_id=3,
            left=100,
            top=300,
            width=600,
            height=200,
            rows=(("A1", "B1"), ("A2", "B2")),
        ),
    ]

    built = walker.walk(shapes, ctx=ctx)

    kinds = [shape.kind for shape in built]
    assert kinds == [ShapeKind.TEXT, ShapeKind.PICTURE, ShapeKind.TABLE]
    assert [shape.z for shape in built] == [0, 1, 2]

    text_shape = built[0]
    assert text_shape.bbox.x_norm == pytest.approx(0.1)
    assert text_shape.text.paras[0].runs[0].text == "Hello"

    picture_shape = built[1]
    assert picture_shape.image.media.data_url == "data:image/png;base64,aW1n"
    assert picture_shape.bbox.w_norm == pytest.approx(0.2)

    table_shape = built[2]
    assert table_shape.table.rows == 2
    cell_texts = [cell.text.paras[0].runs[0].text for cell in table_shape.table.cells]
    assert cell_texts == ["A1", "B1", "A2", "B2"]


def test_slide_walker_handles_group_shapes(ctx: ExtractContext) -> None:
    walker = _make_walker()
    child_text = FakeTextShape(shape_id=10, left=10, top=15, width=200, height=50, text="GroupChild")
    child_picture = FakePictureShape(shape_id=11, left=220, top=25, width=150, height=150)
    group = FakeGroupShape(
        shape_id=9,
        left=200,
        top=100,
        width=400,
        height=300,
        children=[child_text, child_picture],
        name="Group",
    )

    built = walker.walk([group], ctx=ctx)

    assert [shape.kind for shape in built] == [ShapeKind.GROUP, ShapeKind.TEXT, ShapeKind.PICTURE]

    group_shape = built[0]
    assert group_shape.z == 0
    assert group_shape.children == ("s10", "s11")

    text_shape = built[1]
    picture_shape = built[2]
    assert text_shape.group == group_shape.id
    assert picture_shape.group == group_shape.id
    assert text_shape.z == 0
    assert picture_shape.z == 1
    assert text_shape.bbox.x_norm == pytest.approx(0.21)
    assert picture_shape.bbox.y_norm == pytest.approx(0.125)


class FakeFont:
    size = None
    name = None
    bold = None
    italic = None
    underline = None
    color = None


class FakeRun:
    def __init__(self, text: str) -> None:
        self.text = text
        self.font = FakeFont()


class FakeParagraph:
    def __init__(self, text: str) -> None:
        self.runs = [FakeRun(text)]
        self.level = 0
        self.alignment = None


class FakeTextFrame:
    def __init__(self, text: str) -> None:
        self.paragraphs = [FakeParagraph(text)]


class FakeShapeBase:
    def __init__(
        self,
        *,
        shape_id: int,
        left: int,
        top: int,
        width: int,
        height: int,
        name: str | None = None,
    ) -> None:
        self.shape_id = shape_id
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.name = name
        self.rotation = 0


class FakeTextShape(FakeShapeBase):
    shape_type = None

    def __init__(self, *, text: str, **kwargs: int | str | None) -> None:
        super().__init__(**kwargs)
        self.has_text_frame = True
        self.text_frame = FakeTextFrame(text)


class FakeImage:
    def __init__(self) -> None:
        self.blob = b"img"
        self.content_type = "image/png"


class FakePictureShape(FakeShapeBase):
    shape_type = MSO_SHAPE_TYPE.PICTURE

    def __init__(self, **kwargs: int | str | None) -> None:
        super().__init__(**kwargs)
        self.image = FakeImage()
        self.alternative_text = None


class FakeTc:
    def __init__(
        self,
        *,
        grid_span: int | None = None,
        row_span: int | None = None,
        h_merge: bool = False,
        v_merge: bool = False,
    ) -> None:
        self._data = {
            "gridSpan": str(grid_span) if grid_span else None,
            "rowSpan": str(row_span) if row_span else None,
            "hMerge": h_merge,
            "vMerge": v_merge,
        }

    def get(self, key: str) -> str | bool | None:
        return self._data.get(key)


class FakeTableCell:
    def __init__(self, text: str) -> None:
        self._tc = FakeTc()
        self.text_frame = FakeTextFrame(text)


class FakeTable:
    def __init__(self, rows: tuple[tuple[str, ...], ...]) -> None:
        self._cells = [
            [FakeTableCell(text) for text in row]
            for row in rows
        ]
        self.rows = list(range(len(rows)))
        self.columns = list(range(len(rows[0]) if rows else 0))

    def cell(self, r: int, c: int) -> FakeTableCell:
        return self._cells[r][c]


class FakeTableShape(FakeShapeBase):
    shape_type = MSO_SHAPE_TYPE.TABLE
    has_table = True

    def __init__(self, *, rows: tuple[tuple[str, ...], ...], **kwargs: int | str | None) -> None:
        super().__init__(**kwargs)
        self.table = FakeTable(rows)


class FakeGroupShape(FakeShapeBase):
    shape_type = MSO_SHAPE_TYPE.GROUP

    def __init__(self, *, children: list[FakeShapeBase], **kwargs: int | str | None) -> None:
        super().__init__(**kwargs)
        self.shapes = children
