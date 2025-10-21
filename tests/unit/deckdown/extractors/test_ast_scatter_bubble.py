from __future__ import annotations

from pathlib import Path

from deckdown.extractors.ast import AstExtractor
from deckdown.loader import Loader


def _make_xy_bubble(tmp: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import XyChartData, BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE

    p = tmp / "xyb.pptx"
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    xyd = XyChartData()
    xs = [1, 2, 3]
    ys = [4, 5, 6]
    ser = xyd.add_series("S")
    for x, y in zip(xs, ys):
        ser.add_data_point(x, y)
    s1.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1), Inches(4), Inches(3), xyd)

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    bd = BubbleChartData()
    xs = [1, 2, 3]
    ys = [3, 2, 1]
    sz = [10, 20, 30]
    bser = bd.add_series("B")
    for x, y, r in zip(xs, ys, sz):
        bser.add_data_point(x, y, r)
    s2.shapes.add_chart(XL_CHART_TYPE.BUBBLE, Inches(1), Inches(1), Inches(4), Inches(3), bd)

    prs.save(str(p))
    return p


def test_extract_xy_bubble(tmp_path: Path) -> None:
    p = _make_xy_bubble(tmp_path)
    prs = Loader(str(p)).presentation()
    docs = AstExtractor().extract(prs)
    s1 = docs[1].slide
    ch1 = [sh for sh in s1.shapes if getattr(sh, "kind", None).value == "chart"][0]
    assert ch1.chart.type in ("scatter", "xy_scatter", "unknown")
    ser = ch1.chart.series[0]
    assert ser.values and len(ser.values) == 3

    s2 = docs[2].slide
    ch2 = [sh for sh in s2.shapes if getattr(sh, "kind", None).value == "chart"][0]
    assert ch2.chart.type in ("bubble", "unknown")
    ser2 = ch2.chart.series[0]
    assert ser2.values and len(ser2.values) == 3
