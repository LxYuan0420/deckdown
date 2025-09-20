from __future__ import annotations

import sys
from collections.abc import Sequence
from pathlib import Path
from typing import Any


def _require_python_pptx() -> tuple[Any, Any, Any, Any]:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.chart.data import CategoryChartData  # type: ignore
        from pptx.enum.chart import XL_CHART_TYPE  # type: ignore
        from pptx.util import Inches  # type: ignore
    except Exception as exc:  # pragma: no cover - exercised manually
        print(
            "Missing dependency: python-pptx. Install with:\n"
            "  uv pip install python-pptx\n",
            file=sys.stderr,
        )
        raise SystemExit(2) from exc
    return Presentation, Inches, CategoryChartData, XL_CHART_TYPE


def _ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def gen_text_basic(base: Path) -> None:
    presentation, inches, *_ = _require_python_pptx()
    out = base / "text_basic" / "text_basic.pptx"
    _ensure_dir(out.parent)

    prs = presentation()

    # Slide 1: Title and Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Intro"
    body = slide.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = "Welcome to DeckDown"
    p = body.add_paragraph()
    p.text = "Goals: Clean Markdown from PPTX"
    # Bullets
    p = body.add_paragraph()
    p.text = "Item 1"
    p.level = 0
    p = body.add_paragraph()
    p.text = "Sub 1"
    p.level = 1

    # Slide 2: Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Details"
    body = slide.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = "Alpha"
    p = body.add_paragraph()
    p.text = "Beta"

    # Slide 3: Title + Content but no title text (untitled)
    slide = prs.slides.add_slide(slide_layout)
    # leave title empty
    body = slide.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = "Closing"

    prs.save(str(out))


def gen_tables_basic(base: Path) -> None:
    presentation, inches, *_ = _require_python_pptx()
    out = base / "tables_basic" / "tables_basic.pptx"
    _ensure_dir(out.parent)
    prs = presentation()

    # Slide 1: table with escaped cell content
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Tables"
    left, top, width, height = (
        inches(1.0),
        inches(2.0),
        inches(8.0),
        inches(1.5),
    )
    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "A|A"
    table.cell(1, 1).text = "B"
    table.cell(2, 0).text = "C"
    table.cell(2, 1).text = "D"

    # Slide 2: ragged row (leave missing cells empty)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Ragged"
    left, top, width, height = (
        inches(1.0),
        inches(2.0),
        inches(8.0),
        inches(1.5),
    )
    rows, cols = 2, 3
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(0, 2).text = "H3"
    table.cell(1, 0).text = "A"
    # leave (1,1) and (1,2) blank

    prs.save(str(out))


def gen_charts_basic(base: Path) -> None:
    presentation, inches, chart_data, xl_chart_type = _require_python_pptx()
    out = base / "charts_basic" / "charts_basic.pptx"
    _ensure_dir(out.parent)
    prs = presentation()

    # Pie
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Pie"
    data = chart_data()
    data.categories = ["A", "B", "C"]
    data.add_series("Series 1", (30, 20, 50))
    slide.shapes.add_chart(
        xl_chart_type.PIE,
        inches(1),
        inches(2),
        inches(6),
        inches(4),
        data,
    )

    # Column
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Column"
    data = chart_data()
    data.categories = ["Q1", "Q2"]
    data.add_series("S1", (10, 12))
    data.add_series("S2", (8, 9))
    slide.shapes.add_chart(
        xl_chart_type.COLUMN_CLUSTERED,
        inches(1),
        inches(2),
        inches(6),
        inches(4),
        data,
    )

    # Line
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Line"
    data = chart_data()
    data.categories = ["Jan", "Feb", "Mar"]
    data.add_series("S", (1, 2, 3))
    slide.shapes.add_chart(
        xl_chart_type.LINE,
        inches(1),
        inches(2),
        inches(6),
        inches(4),
        data,
    )

    prs.save(str(out))


def gen_diff_pair(base: Path) -> None:
    presentation, inches, *_ = _require_python_pptx()
    v1 = base / "diff_pair" / "v1" / "diff_pair_v1.pptx"
    v2 = base / "diff_pair" / "v2" / "diff_pair_v2.pptx"
    _ensure_dir(v1.parent)
    _ensure_dir(v2.parent)

    prs = presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Intro"
    body = slide.placeholders[1].text_frame
    body.text = "Welcome"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data"
    body = slide.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = "Item 1"
    p.level = 0
    p = body.add_paragraph()
    p.text = "Item 2"
    p.level = 0
    prs.save(str(v1))

    prs2 = presentation(str(v1))
    s1 = prs2.slides[0]
    s1.placeholders[1].text_frame.paragraphs[0].text = "Welcome!"
    s2 = prs2.slides[1]
    tf = s2.placeholders[1].text_frame
    tf.paragraphs[1].text = "Item 2 (updated)"
    p = tf.add_paragraph()
    p.text = "Item 3"
    p.level = 0
    prs2.save(str(v2))


def gen_notes(base: Path) -> None:
    presentation, inches, *_ = _require_python_pptx()
    out = base / "notes" / "notes.pptx"
    _ensure_dir(out.parent)
    prs = presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Notes"
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Confidential: internal only"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "More"
    slide.notes_slide.notes_text_frame.text = "Next steps: ship M1"
    prs.save(str(out))


def gen_edge_cases(base: Path) -> None:
    presentation, inches, *_ = _require_python_pptx()
    out = base / "edge_cases" / "edge_cases.pptx"
    _ensure_dir(out.parent)
    prs = presentation()

    # Slide 1: no title, no shapes (blank layout)
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Slide 2: special chars
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Special Chars"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = "Pipe | Asterisk * Underscore _ Backtick `"
    p = tf.add_paragraph()
    p.text = "Line1"
    p = tf.add_paragraph()
    p.text = "Line2"

    # Slide 3: deep bullets
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Deep Bullets"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for lvl, text in enumerate(["L0", "L1", "L2", "L3"]):
        p = tf.add_paragraph() if lvl else tf.paragraphs[0]
        p.text = text
        p.level = lvl

    # Slide 4: table with newline and pipe
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Table Weird"
    left, top, width, height = (
        inches(1.0),
        inches(2.0),
        inches(8.0),
        inches(1.5),
    )
    rows, cols = 2, 2
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "A|A"
    table.cell(1, 1).text = "B\nB"

    prs.save(str(out))


def main(argv: Sequence[str] | None = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(description="Generate sample PPTX decks under data/samples/")
    parser.add_argument("--out", default="data/samples", help="Base output directory")
    parser.add_argument(
        "--only",
        default=[],
        nargs="*",
        choices=[
            "text_basic",
            "tables_basic",
            "charts_basic",
            "diff_pair",
            "notes",
            "edge_cases",
        ],
        help="Subset of scenarios to generate (default: all)",
    )
    ns = parser.parse_args(list(argv) if argv is not None else None)

    base = Path(ns.out)
    scenarios = ns.only or [
        "text_basic",
        "tables_basic",
        "charts_basic",
        "diff_pair",
        "notes",
        "edge_cases",
    ]

    for sc in scenarios:
        if sc == "text_basic":
            gen_text_basic(base)
        elif sc == "tables_basic":
            gen_tables_basic(base)
        elif sc == "charts_basic":
            gen_charts_basic(base)
        elif sc == "diff_pair":
            gen_diff_pair(base)
        elif sc == "notes":
            gen_notes(base)
        elif sc == "edge_cases":
            gen_edge_cases(base)

    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
