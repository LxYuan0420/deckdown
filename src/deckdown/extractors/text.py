from __future__ import annotations

from collections.abc import Iterable
from typing import Any

from deckdown.models import Bullet, Deck, Slide, TextBlock


def extract_deck_text(path: str, *, with_notes: bool = False) -> Deck:
    """Extract titles, text paragraphs, and bullets into a Deck.

    Heuristics:
    - Reading order: shapes are processed top-to-bottom, then left-to-right.
    - Title: derived from slide.shapes.title when available.
    - Text vs bullets: for each text frame, if a paragraph with level>0 exists,
      we treat the preceding paragraph (if any) as the level-0 bullet start and
      everything from there as bullets; earlier paragraphs are aggregated as a
      single TextBlock. If no paragraph has level>0, all lines are text.
    """
    prs = _load_presentation(path)
    slides_out = [
        _extract_slide(i, s) for i, s in enumerate(prs.slides, start=1)
    ]
    return Deck(file=str(path), title=_presentation_title(prs) or None, slides=tuple(slides_out))


def _extract_slide(index: int, slide: Any) -> Slide:  # noqa: ANN401
    title = _get_slide_title(slide)
    text_blocks: list[TextBlock] = []
    bullets: list[Bullet] = []

    title_shape = getattr(slide.shapes, "title", None)

    for shp in _iter_text_shapes_in_order(slide):
        if title_shape is not None and shp is title_shape:
            continue

        tf = shp.text_frame
        paras: list[tuple[str, int]] = []
        for p in tf.paragraphs:
            txt = (p.text or "").strip()
            if txt == "":
                continue
            # paragraph.level is 0-based
            level = int(getattr(p, "level", 0) or 0)
            paras.append((txt, level))

        if not paras:
            continue

        # Skip if this shape appears to be just the slide title repeated
        if title and len(paras) == 1 and paras[0][0] == title:
            continue

        first_lvl_gt0 = next((i for i, (_, lvl) in enumerate(paras) if lvl > 0), None)
        if first_lvl_gt0 is None:
            # All are text; aggregate into a single TextBlock with newlines
            text_blocks.append(TextBlock(text="\n".join(t for t, _ in paras)))
        else:
            start_idx = max(first_lvl_gt0 - 1, 0)
            if start_idx > 0:
                text_blocks.append(TextBlock(text="\n".join(t for t, _ in paras[:start_idx])))
            for t, lvl in paras[start_idx:]:
                bullets.append(Bullet(level=lvl, text=t))

    return Slide(
        index=index,
        title=title,
        text_blocks=tuple(text_blocks),
        bullets=tuple(bullets),
        tables=(),
        charts=(),
    )


# Internal helpers using python-pptx (imported lazily for testability)

def _load_presentation(path: str) -> Any:  # noqa: ANN401
    from pptx import Presentation

    return Presentation(path)


def _presentation_title(prs: Any) -> str | None:  # noqa: ANN401
    # Prefer the first slide title if present; otherwise, basename without extension
    if prs.slides:
        t = _get_slide_title(prs.slides[0])
        if t:
            return t
    try:
        # Presentation has a 'core_properties' title sometimes set
        title = prs.core_properties.title
        if title:
            return str(title)
    except Exception:
        return None
    return None


def _get_slide_title(slide: Any) -> str | None:  # noqa: ANN401
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        text = getattr(title_shape.text_frame.paragraphs[0], "text", "")
        t = (text or "").strip()
        return t or None
    return None


def _iter_text_shapes_in_order(slide: Any) -> Iterable[Any]:  # noqa: ANN401
    shapes = [
        shp
        for shp in slide.shapes
        if getattr(shp, "has_text_frame", False)
    ]
    shapes.sort(key=lambda s: (int(getattr(s, "top", 0)), int(getattr(s, "left", 0))))
    return shapes
