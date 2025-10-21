from __future__ import annotations

from typing import Any, Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckdown.ast import BasicShape, BasicStyle, FillSpec, LineShape, ShapeKind, StrokeSpec, TextPayload
from deckdown.extractors.context import ExtractContext
from deckdown.extractors.handlers.base import ShapeHandler
from deckdown.extractors.utils import extract_text_payload


def _color_from_fill(fill: Any, ctx: ExtractContext) -> dict | None:  # noqa: ANN401
    try:
        fc = getattr(fill, "fore_color", None)
        if fc is not None:
            data = ctx.theme.color_dict_from_colorformat(fc)
            if data:
                return data
    except Exception:
        return None
    return None


def _basic_style(shp: Any, ctx: ExtractContext) -> BasicStyle | None:  # noqa: ANN401
    try:
        fill_color = _color_from_fill(getattr(shp, "fill", None), ctx)
        line = getattr(shp, "line", None)
        stroke_color = None
        width_pt = None
        dash = None
        if line is not None:
            try:
                stroke_color = _color_from_fill(getattr(line, "fill", None), ctx)
            except Exception:
                pass
            try:
                if getattr(line, "width", None) is not None:
                    width_pt = round(float(line.width.pt), 2)  # type: ignore[union-attr]
            except Exception:
                pass
        if fill_color or stroke_color or width_pt:
            return BasicStyle(
                fill=FillSpec(color=fill_color) if fill_color else None,
                stroke=StrokeSpec(color=stroke_color, width_pt=width_pt, dash=dash),
            )
    except Exception:
        return None
    return None


class BasicShapeHandler(ShapeHandler):
    def supports(self, shape: Any) -> bool:  # noqa: ANN401
        return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE

    def build(self, shape: Any, *, z: int, ctx: ExtractContext) -> Optional[BasicShape]:  # noqa: ANN401
        bbox = ctx.bbox(
            left_emu=int(getattr(shape, "left", 0)),
            top_emu=int(getattr(shape, "top", 0)),
            width_emu=int(getattr(shape, "width", 0)),
            height_emu=int(getattr(shape, "height", 0)),
        )
        geom = str(getattr(getattr(shape, "auto_shape_type", None), "name", None) or "autoShape").lower()
        style = _basic_style(shape, ctx)
        text_payload: TextPayload | None = None
        if getattr(shape, "has_text_frame", False):
            text_payload = extract_text_payload(shape.text_frame, ctx.theme)
        rot = None
        try:
            rot = float(getattr(shape, "rotation"))  # type: ignore[arg-type]
        except Exception:
            rot = None
        return BasicShape(
            id=f"s{getattr(shape, 'shape_id', z)}",
            kind=ShapeKind.BASIC,
            name=getattr(shape, "name", None),
            bbox=bbox,
            z=z,
            rotation=rot,
            geom=geom,
            style=style,
            text=text_payload,
        )


class LineShapeHandler(ShapeHandler):
    def supports(self, shape: Any) -> bool:  # noqa: ANN401
        return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE

    def build(self, shape: Any, *, z: int, ctx: ExtractContext) -> Optional[LineShape]:  # noqa: ANN401
        bbox = ctx.bbox(
            left_emu=int(getattr(shape, "left", 0)),
            top_emu=int(getattr(shape, "top", 0)),
            width_emu=int(getattr(shape, "width", 0)),
            height_emu=int(getattr(shape, "height", 0)),
        )
        style = _basic_style(shape, ctx)
        rot = None
        try:
            rot = float(getattr(shape, "rotation"))  # type: ignore[arg-type]
        except Exception:
            rot = None
        return LineShape(
            id=f"s{getattr(shape, 'shape_id', z)}",
            kind=ShapeKind.LINE,
            name=getattr(shape, "name", None),
            bbox=bbox,
            z=z,
            rotation=rot,
            style=style,
        )
