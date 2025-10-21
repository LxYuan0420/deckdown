from __future__ import annotations

from typing import Any, Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckdown.ast import ShapeKind, TextPayload, TextShape
from deckdown.extractors.context import ExtractContext
from deckdown.extractors.handlers.base import ShapeHandler
from deckdown.extractors.utils import extract_text_payload


class TextShapeHandler(ShapeHandler):
    def supports(self, shape: Any) -> bool:  # noqa: ANN401
        st = getattr(shape, "shape_type", None)
        if st in (MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.AUTO_SHAPE):
            return False
        return bool(getattr(shape, "has_text_frame", False))

    def build(self, shape: Any, *, z: int, ctx: ExtractContext) -> Optional[TextShape]:  # noqa: ANN401
        bbox = ctx.bbox(
            left_emu=int(getattr(shape, "left", 0)),
            top_emu=int(getattr(shape, "top", 0)),
            width_emu=int(getattr(shape, "width", 0)),
            height_emu=int(getattr(shape, "height", 0)),
        )
        text: TextPayload = extract_text_payload(shape.text_frame, ctx.theme)
        rot = None
        try:
            rot = float(getattr(shape, "rotation"))  # type: ignore[arg-type]
        except Exception:
            rot = None
        return TextShape(
            id=f"s{getattr(shape, 'shape_id', z)}",
            kind=ShapeKind.TEXT,
            name=getattr(shape, "name", None),
            bbox=bbox,
            z=z,
            rotation=rot,
            text=text,
        )
