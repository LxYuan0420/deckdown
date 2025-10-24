from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Iterable

from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckdown.ast import Shape, SlideDoc, SlideModel, SlideSize
from deckdown.extractors.context import ExtractContext
from deckdown.extractors.handlers.base import ShapeHandler
from deckdown.extractors.handlers.table_handler import TableShapeHandler
from deckdown.extractors.handlers.chart_handler import ChartShapeHandler
from deckdown.extractors.handlers.picture_handler import PictureShapeHandler
from deckdown.extractors.handlers.basic_line_handler import BasicShapeHandler, LineShapeHandler
from deckdown.extractors.handlers.text_handler import TextShapeHandler
from deckdown.extractors.group import GroupExtractor
from deckdown.color.theme import ThemeResolver
from deckdown.media import AssetStore, MediaEmbedMode


@dataclass(frozen=True)
class AstExtractor:
    """Build per-slide AST using focused shape handlers.

    Single path implementation (no legacy fallback). Handlers encapsulate
    per-shape logic to keep this extractor small and readable.
    """

    media_mode: MediaEmbedMode = "base64"
    asset_store: AssetStore | None = None

    def extract(self, prs: Any) -> dict[int, SlideDoc]:  # noqa: ANN401
        size = SlideSize(width_emu=int(prs.slide_width), height_emu=int(prs.slide_height))
        ctx = ExtractContext(
            size=size,
            theme=ThemeResolver.from_presentation(prs),
            media_mode=self.media_mode,
            asset_store=self.asset_store,
        )

        handlers: tuple[ShapeHandler, ...] = (
            TableShapeHandler(),
            ChartShapeHandler(),
            PictureShapeHandler(),
            LineShapeHandler(),
            BasicShapeHandler(),
            TextShapeHandler(),
        )
        group_extractor = GroupExtractor(handlers=handlers)
        walker = SlideWalker(handlers=handlers, group_extractor=group_extractor)

        out: dict[int, SlideDoc] = {}
        for idx, slide in enumerate(prs.slides, start=1):
            walked = walker.walk(slide.shapes, ctx=ctx)
            out[idx] = SlideDoc(slide=SlideModel(index=idx, size=size, shapes=tuple(walked)))
        return out


@dataclass(frozen=True)
class SlideWalker:
    handlers: tuple[ShapeHandler, ...]
    group_extractor: GroupExtractor

    def walk(self, shapes: Iterable[Any], *, ctx: ExtractContext) -> list[Shape]:  # noqa: ANN401
        built: list[Shape] = []
        z = 0
        for shp in shapes:
            z = self._dispatch(shp, current_z=z, ctx=ctx, out=built)
        return built

    def _dispatch(
        self,
        shp: Any,  # noqa: ANN401
        *,
        current_z: int,
        ctx: ExtractContext,
        out: list[Shape],
    ) -> int:
        if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            children, next_z, group_shape = self.group_extractor.extract(shp, z_start=current_z, ctx=ctx)
            out.append(group_shape)
            out.extend(children)
            return next_z

        for handler in self.handlers:
            if not handler.supports(shp):
                continue
            built = handler.build(shp, z=current_z, ctx=ctx)
            if built is not None:
                out.append(built)
            return current_z + 1

        return current_z
