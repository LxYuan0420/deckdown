from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.util import Emu

from deckdown.ast import BasicShape, LineShape, PictureShape, SlideDoc, TableShape, TextShape, ChartShape
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


@dataclass(frozen=True)
class DeckAssembler:
    def assemble(self, docs: Iterable[SlideDoc], *, out: Path) -> None:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        # Ensure empty deck
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId  # type: ignore[attr-defined]
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])  # type: ignore[attr-defined]

        # Set deck slide size from first doc if present
        docs_list = list(docs)
        if docs_list:
            first = docs_list[0]
            try:
                prs.slide_width = Emu(first.slide.size.width_emu)
                prs.slide_height = Emu(first.slide.size.height_emu)
            except Exception:
                pass

        for doc in docs_list:
            s = prs.slides.add_slide(blank)
            # Note: slide size is a deck-level setting in PPTX; we keep default for now.
            for sh in doc.slide.shapes:
                if isinstance(sh, TextShape):
                    self._add_text(s, sh)
                elif isinstance(sh, PictureShape):
                    self._add_picture(s, sh)
                elif isinstance(sh, TableShape):
                    self._add_table(s, sh)
                elif isinstance(sh, BasicShape):
                    self._add_basic(s, sh)
                elif isinstance(sh, LineShape):
                    self._add_line(s, sh)
                elif isinstance(sh, ChartShape):
                    self._add_chart(s, sh)

        prs.save(str(out))

    # --- helpers ---
    def _add_text(self, slide, sh: TextShape) -> None:  # noqa: ANN001
        left = Emu(sh.bbox.x_emu); top = Emu(sh.bbox.y_emu); width = Emu(sh.bbox.w_emu); height = Emu(sh.bbox.h_emu)
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.clear()
        first = True
        for p in sh.text.paras:
            if first:
                par = tf.paragraphs[0]
                first = False
            else:
                par = tf.add_paragraph()
            par.level = p.lvl
            par.text = "".join(r.text for r in p.runs)

    def _add_picture(self, slide, sh: PictureShape) -> None:  # noqa: ANN001
        media = sh.image.media
        if not media or not media.data_url:
            return
        data_url = media.data_url
        if not data_url.startswith("data:") or ";base64," not in data_url:
            return
        b64 = data_url.split(",", 1)[1]
        try:
            data = BytesIO(__import__("base64").b64decode(b64))
            data.seek(0)
        except Exception:
            return
        left = Emu(sh.bbox.x_emu); top = Emu(sh.bbox.y_emu)
        width = Emu(sh.bbox.w_emu); height = Emu(sh.bbox.h_emu)
        slide.shapes.add_picture(data, left, top, width=width, height=height)

    def _add_table(self, slide, sh: TableShape) -> None:  # noqa: ANN001
        left = Emu(sh.bbox.x_emu); top = Emu(sh.bbox.y_emu); width = Emu(sh.bbox.w_emu); height = Emu(sh.bbox.h_emu)
        rows, cols = sh.table.rows, sh.table.cols
        t = slide.shapes.add_table(rows, cols, left, top, width, height).table
        # Initialize with empty strings
        for r in range(rows):
            for c in range(cols):
                t.cell(r, c).text = ""
        # Apply merges and set text on top-left cells
        for cell in sh.table.cells:
            r, c = cell.r, cell.c
            rr = r + max(1, cell.rowspan) - 1
            cc = c + max(1, cell.colspan) - 1
            if rr > r or cc > c:
                try:
                    t.cell(r, c).merge(t.cell(rr, cc))
                except Exception:
                    pass
            # write text
            tf = t.cell(r, c).text_frame
            tf.clear()
            first = True
            for p in cell.text.paras:
                if first:
                    par = tf.paragraphs[0]
                    first = False
                else:
                    par = tf.add_paragraph()
                par.level = p.lvl
                par.text = "".join(run.text for run in p.runs)

    def _add_basic(self, slide, sh: BasicShape) -> None:  # noqa: ANN001
        left = Emu(sh.bbox.x_emu); top = Emu(sh.bbox.y_emu); width = Emu(sh.bbox.w_emu); height = Emu(sh.bbox.h_emu)
        geom_map = {
            "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            "rounded rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            "roundrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
            "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
        }
        kind = geom_map.get((sh.geom or "").lower(), MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        shp = slide.shapes.add_shape(kind, left, top, width, height)
        # text
        if sh.text and sh.text.paras:
            tf = shp.text_frame
            tf.clear()
            first = True
            for p in sh.text.paras:
                if first:
                    par = tf.paragraphs[0]
                    first = False
                else:
                    par = tf.add_paragraph()
                par.level = p.lvl
                par.text = "".join(run.text for run in p.runs)
        # style (fill only; stroke best-effort)
        try:
            if sh.style and sh.style.fill and sh.style.fill.color and sh.style.fill.color.get("resolved_rgb"):
                shp.fill.solid()
                from pptx.dml.color import RGBColor

                rgb = sh.style.fill.color["resolved_rgb"][1:]
                shp.fill.fore_color.rgb = RGBColor(int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16))
        except Exception:
            pass

        # Per-point override colors and axes
        try:
            plot0 = chart.plots[0]
            from pptx.dml.color import RGBColor
            for i, ser in enumerate((sh.chart.series or ())) :
                chart_ser = plot0.series[i]
                if getattr(ser, "points", None):
                    for ptmeta in ser.points:
                        idx = ptmeta.get("idx")
                        col = (ptmeta.get("color") or {}).get("resolved_rgb", None)
                        if idx is None or not col:
                            continue
                        hexrgb = col[1:]
                        try:
                            pt = chart_ser.points[idx]
                            pt.format.fill.solid()
                            pt.format.fill.fore_color.rgb = RGBColor(
                                int(hexrgb[0:2], 16), int(hexrgb[2:4], 16), int(hexrgb[4:6], 16)
                            )
                        except Exception:
                            continue
            # Axes (titles, ranges, number format)
            ax = sh.chart.axes or {}
            if "category" in ax:
                t = ax["category"].get("title")
                if t:
                    chart.category_axis.has_title = True
                    chart.category_axis.axis_title.text_frame.text = str(t)
            if "value" in ax:
                v = ax["value"]
                if v.get("title"):
                    chart.value_axis.has_title = True
                    chart.value_axis.axis_title.text_frame.text = str(v["title"])
                if v.get("min") is not None:
                    chart.value_axis.minimum_scale = float(v["min"])
                if v.get("max") is not None:
                    chart.value_axis.maximum_scale = float(v["max"])
                if v.get("major_unit") is not None:
                    chart.value_axis.major_unit = float(v["major_unit"])
                if v.get("format_code"):
                    chart.value_axis.tick_labels.number_format = str(v["format_code"])
        except Exception:
            pass

        # Per-point override colors
        try:
            plot0 = chart.plots[0]
            from pptx.dml.color import RGBColor
            for i, ser in enumerate((sh.chart.series or ())) :
                if not getattr(ser, "points", None):
                    continue
                chart_ser = plot0.series[i]
                for ptmeta in ser.points:
                    idx = ptmeta.get("idx")
                    col = (ptmeta.get("color") or {}).get("resolved_rgb", None)
                    if idx is None or not col:
                        continue
                    hexrgb = col[1:]
                    try:
                        pt = chart_ser.points[idx]
                        pt.format.fill.solid()
                        pt.format.fill.fore_color.rgb = RGBColor(
                            int(hexrgb[0:2], 16), int(hexrgb[2:4], 16), int(hexrgb[4:6], 16)
                        )
                    except Exception:
                        continue
        except Exception:
            pass
    def _add_line(self, slide, sh: LineShape) -> None:  # noqa: ANN001
        x1 = Emu(sh.bbox.x_emu); y1 = Emu(sh.bbox.y_emu)
        x2 = Emu(sh.bbox.x_emu + sh.bbox.w_emu); y2 = Emu(sh.bbox.y_emu + sh.bbox.h_emu)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)

    def _add_chart(self, slide, sh: ChartShape) -> None:  # noqa: ANN001
        ct = (sh.chart.type or "unknown").lower()
        ctype_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "donut": XL_CHART_TYPE.DOUGHNUT,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "bubble": XL_CHART_TYPE.BUBBLE,
        }
        xl_type = ctype_map.get(ct, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Build data by chart family
        if xl_type in (XL_CHART_TYPE.XY_SCATTER,):
            data = XyChartData()
            for ser in sh.chart.series or ():
                s = data.add_series(ser.name or "Series")
                xs = list(ser.x_values or ())
                ys = list(ser.values or ())
                n = min(len(xs), len(ys))
                for i in range(n):
                    xv = xs[i] if xs[i] is not None else 0
                    yv = ys[i] if ys[i] is not None else 0
                    s.add_data_point(xv, yv)
        elif xl_type in (XL_CHART_TYPE.BUBBLE,):
            data = BubbleChartData()
            for ser in sh.chart.series or ():
                s = data.add_series(ser.name or "Series")
                xs = list(ser.x_values or ())
                ys = list(ser.values or ())
                sz = list(ser.sizes or ())
                n = min(len(xs), len(ys), len(sz) if sz else len(ys))
                for i in range(n):
                    xv = xs[i] if xs[i] is not None else 0
                    yv = ys[i] if ys[i] is not None else 0
                    rv = sz[i] if (sz and sz[i] is not None) else 1
                    s.add_data_point(xv, yv, rv)
        else:
            data = CategoryChartData()
            data.categories = list(sh.chart.categories or [])
            for ser in sh.chart.series or ():
                name = ser.name or "Series"
                values = [v if v is not None else 0 for v in (ser.values or ())]
                data.add_series(name, values)

        left = Emu(sh.bbox.x_emu); top = Emu(sh.bbox.y_emu); width = Emu(sh.bbox.w_emu); height = Emu(sh.bbox.h_emu)
        chart = slide.shapes.add_chart(xl_type, left, top, width, height, data).chart

        # Legend and data labels
        if sh.chart.plot_area:
            has_labels = bool(sh.chart.plot_area.get("has_data_labels", False))
            if chart.plots:
                try:
                    chart.plots[0].has_data_labels = has_labels
                except Exception:
                    pass
            has_legend = bool(sh.chart.plot_area.get("has_legend", False))
            chart.has_legend = has_legend
            pos = (sh.chart.plot_area.get("legend_pos") or "").lower()
            pos_map = {
                "right": XL_LEGEND_POSITION.RIGHT,
                "left": XL_LEGEND_POSITION.LEFT,
                "top": XL_LEGEND_POSITION.TOP,
                "bottom": XL_LEGEND_POSITION.BOTTOM,
            }
            if has_legend and pos in pos_map:
                try:
                    chart.legend.position = pos_map[pos]
                except Exception:
                    pass

        # Series data labels and colors (explicit)
        try:
            plot0 = chart.plots[0]
            from pptx.dml.color import RGBColor
            for i, ser in enumerate((sh.chart.series or ())) :
                chart_ser = plot0.series[i]
                lab = getattr(ser, "labels", None) or {}
                if lab:
                    dl = chart_ser.data_labels
                    for key in ("show_value", "show_category_name", "show_series_name", "show_percentage"):
                        if key in lab and lab[key] is not None:
                            try:
                                setattr(dl, key, bool(lab[key]))
                            except Exception:
                                pass
                    if lab.get("number_format"):
                        try:
                            dl.number_format = str(lab["number_format"])
                        except Exception:
                            pass
                if ser.color and ser.color.get("resolved_rgb"):
                    rgb = ser.color["resolved_rgb"][1:]
                    chart_ser.format.fill.solid()
                    chart_ser.format.fill.fore_color.rgb = RGBColor(
                        int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16)
                    )
        except Exception:
            pass
